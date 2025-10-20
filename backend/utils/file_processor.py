"""
File processing utilities for document uploads and text extraction.

This module provides comprehensive file processing capabilities for the Personal Learning Agent,
supporting multiple file formats and text extraction for skills assessment artifacts.
"""

import os
import io
import mimetypes
import hashlib
from typing import Dict, List, Optional, Union, Any, Tuple
from dataclasses import dataclass
from pathlib import Path
import logging

# File processing imports
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl
import pandas as pd
from bs4 import BeautifulSoup
import markdown
import json

# FastAPI imports for file handling
from fastapi import UploadFile, HTTPException

from ..core.config import get_config

# Configure logging
logger = logging.getLogger(__name__)


@dataclass
class FileMetadata:
    """Metadata for processed files."""
    filename: str
    file_size: int
    file_type: str
    mime_type: str
    file_hash: str
    processing_time: float
    text_length: int
    page_count: Optional[int] = None
    error: Optional[str] = None


@dataclass
class ProcessedContent:
    """Processed content from files."""
    text: str
    metadata: FileMetadata
    structured_data: Optional[Dict[str, Any]] = None
    extracted_entities: Optional[List[str]] = None


class FileProcessor:
    """
    Comprehensive file processor for multiple document formats.
    
    Supports PDF, DOCX, PPTX, XLSX, TXT, HTML, MD, and JSON files
    with text extraction and metadata generation.
    """
    
    # Supported file types and their extensions
    SUPPORTED_EXTENSIONS = {
        '.pdf': 'application/pdf',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.doc': 'application/msword',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.ppt': 'application/vnd.ms-powerpoint',
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.xls': 'application/vnd.ms-excel',
        '.txt': 'text/plain',
        '.html': 'text/html',
        '.htm': 'text/html',
        '.md': 'text/markdown',
        '.json': 'application/json',
        '.csv': 'text/csv'
    }
    
    # Maximum file size (10MB)
    MAX_FILE_SIZE = 10 * 1024 * 1024
    
    def __init__(self, config_manager=None):
        """
        Initialize file processor.
        
        Args:
            config_manager: Configuration manager instance
        """
        self.config = config_manager or get_config()
        self.settings = self.config.get_settings()
        
        # Create upload directory if it doesn't exist
        self.upload_dir = Path(self.settings.data_dir) / "uploads"
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        
        logger.info(f"File processor initialized with upload directory: {self.upload_dir}")
    
    def validate_file(self, file: UploadFile) -> Tuple[bool, str]:
        """
        Validate uploaded file.
        
        Args:
            file: FastAPI UploadFile object
            
        Returns:
            Tuple of (is_valid, error_message)
        """
        try:
            # Check file size
            if hasattr(file, 'size') and file.size and file.size > self.MAX_FILE_SIZE:
                return False, f"File size exceeds maximum limit of {self.MAX_FILE_SIZE / (1024*1024):.1f}MB"
            
            # Check file extension
            if not file.filename:
                return False, "No filename provided"
            
            file_ext = Path(file.filename).suffix.lower()
            if file_ext not in self.SUPPORTED_EXTENSIONS:
                return False, f"Unsupported file type: {file_ext}. Supported types: {', '.join(self.SUPPORTED_EXTENSIONS.keys())}"
            
            # Check MIME type if available
            if hasattr(file, 'content_type') and file.content_type:
                expected_mime = self.SUPPORTED_EXTENSIONS.get(file_ext)
                if expected_mime and file.content_type != expected_mime:
                    logger.warning(f"MIME type mismatch for {file.filename}: expected {expected_mime}, got {file.content_type}")
            
            return True, ""
            
        except Exception as e:
            logger.error(f"File validation error: {e}")
            return False, f"File validation failed: {str(e)}"
    
    def calculate_file_hash(self, content: bytes) -> str:
        """
        Calculate SHA-256 hash of file content.
        
        Args:
            content: File content as bytes
            
        Returns:
            SHA-256 hash string
        """
        return hashlib.sha256(content).hexdigest()
    
    def extract_text_from_pdf(self, content: bytes) -> Tuple[str, int]:
        """
        Extract text from PDF file.
        
        Args:
            content: PDF file content as bytes
            
        Returns:
            Tuple of (extracted_text, page_count)
        """
        try:
            pdf_file = io.BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text = ""
            page_count = len(pdf_reader.pages)
            
            for page_num in range(page_count):
                try:
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text() + "\n"
                except Exception as e:
                    logger.warning(f"Error extracting text from page {page_num + 1}: {e}")
                    continue
            
            return text.strip(), page_count
            
        except Exception as e:
            logger.error(f"PDF text extraction failed: {e}")
            raise Exception(f"PDF processing failed: {str(e)}")
    
    def extract_text_from_docx(self, content: bytes) -> Tuple[str, int]:
        """
        Extract text from DOCX file.
        
        Args:
            content: DOCX file content as bytes
            
        Returns:
            Tuple of (extracted_text, page_count)
        """
        try:
            doc_file = io.BytesIO(content)
            doc = Document(doc_file)
            
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Estimate page count (rough approximation)
            page_count = max(1, len(text) // 2000)  # ~2000 characters per page
            
            return text.strip(), page_count
            
        except Exception as e:
            logger.error(f"DOCX text extraction failed: {e}")
            raise Exception(f"DOCX processing failed: {str(e)}")
    
    def extract_text_from_pptx(self, content: bytes) -> Tuple[str, int]:
        """
        Extract text from PPTX file.
        
        Args:
            content: PPTX file content as bytes
            
        Returns:
            Tuple of (extracted_text, slide_count)
        """
        try:
            pptx_file = io.BytesIO(content)
            presentation = Presentation(pptx_file)
            
            text = ""
            slide_count = len(presentation.slides)
            
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return text.strip(), slide_count
            
        except Exception as e:
            logger.error(f"PPTX text extraction failed: {e}")
            raise Exception(f"PPTX processing failed: {str(e)}")
    
    def extract_text_from_xlsx(self, content: bytes) -> Tuple[str, int]:
        """
        Extract text from XLSX file.
        
        Args:
            content: XLSX file content as bytes
            
        Returns:
            Tuple of (extracted_text, sheet_count)
        """
        try:
            xlsx_file = io.BytesIO(content)
            workbook = openpyxl.load_workbook(xlsx_file, data_only=True)
            
            text = ""
            sheet_count = len(workbook.sheetnames)
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        text += row_text + "\n"
                text += "\n"
            
            return text.strip(), sheet_count
            
        except Exception as e:
            logger.error(f"XLSX text extraction failed: {e}")
            raise Exception(f"XLSX processing failed: {str(e)}")
    
    def extract_text_from_txt(self, content: bytes) -> Tuple[str, int]:
        """
        Extract text from TXT file.
        
        Args:
            content: TXT file content as bytes
            
        Returns:
            Tuple of (extracted_text, line_count)
        """
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            text = ""
            
            for encoding in encodings:
                try:
                    text = content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if not text:
                # Fallback to utf-8 with error handling
                text = content.decode('utf-8', errors='replace')
            
            line_count = len(text.splitlines())
            return text.strip(), line_count
            
        except Exception as e:
            logger.error(f"TXT text extraction failed: {e}")
            raise Exception(f"TXT processing failed: {str(e)}")
    
    def extract_text_from_html(self, content: bytes) -> Tuple[str, int]:
        """
        Extract text from HTML file.
        
        Args:
            content: HTML file content as bytes
            
        Returns:
            Tuple of (extracted_text, element_count)
        """
        try:
            # Decode HTML content
            html_text = content.decode('utf-8', errors='replace')
            
            # Parse HTML and extract text
            soup = BeautifulSoup(html_text, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text content
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            element_count = len(soup.find_all())
            
            return text.strip(), element_count
            
        except Exception as e:
            logger.error(f"HTML text extraction failed: {e}")
            raise Exception(f"HTML processing failed: {str(e)}")
    
    def extract_text_from_markdown(self, content: bytes) -> Tuple[str, int]:
        """
        Extract text from Markdown file.
        
        Args:
            content: Markdown file content as bytes
            
        Returns:
            Tuple of (extracted_text, line_count)
        """
        try:
            # Decode markdown content
            md_text = content.decode('utf-8', errors='replace')
            
            # Convert markdown to HTML then extract text
            html = markdown.markdown(md_text)
            soup = BeautifulSoup(html, 'html.parser')
            
            # Extract text content
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            line_count = len(md_text.splitlines())
            
            return text.strip(), line_count
            
        except Exception as e:
            logger.error(f"Markdown text extraction failed: {e}")
            raise Exception(f"Markdown processing failed: {str(e)}")
    
    def extract_text_from_json(self, content: bytes) -> Tuple[str, int]:
        """
        Extract text from JSON file.
        
        Args:
            content: JSON file content as bytes
            
        Returns:
            Tuple of (extracted_text, object_count)
        """
        try:
            # Decode JSON content
            json_text = content.decode('utf-8', errors='replace')
            
            # Parse JSON
            data = json.loads(json_text)
            
            # Extract text from JSON structure
            def extract_text_from_dict(obj, path=""):
                text_parts = []
                if isinstance(obj, dict):
                    for key, value in obj.items():
                        current_path = f"{path}.{key}" if path else key
                        if isinstance(value, (dict, list)):
                            text_parts.extend(extract_text_from_dict(value, current_path))
                        else:
                            text_parts.append(f"{current_path}: {str(value)}")
                elif isinstance(obj, list):
                    for i, item in enumerate(obj):
                        current_path = f"{path}[{i}]" if path else f"[{i}]"
                        if isinstance(item, (dict, list)):
                            text_parts.extend(extract_text_from_dict(item, current_path))
                        else:
                            text_parts.append(f"{current_path}: {str(item)}")
                else:
                    text_parts.append(f"{path}: {str(obj)}")
                return text_parts
            
            text_parts = extract_text_from_dict(data)
            text = "\n".join(text_parts)
            
            # Count objects (rough approximation)
            object_count = len(str(data).split(','))
            
            return text.strip(), object_count
            
        except Exception as e:
            logger.error(f"JSON text extraction failed: {e}")
            raise Exception(f"JSON processing failed: {str(e)}")
    
    def extract_text_from_csv(self, content: bytes) -> Tuple[str, int]:
        """
        Extract text from CSV file.
        
        Args:
            content: CSV file content as bytes
            
        Returns:
            Tuple of (extracted_text, row_count)
        """
        try:
            # Decode CSV content
            csv_text = content.decode('utf-8', errors='replace')
            
            # Read CSV with pandas
            csv_file = io.StringIO(csv_text)
            df = pd.read_csv(csv_file)
            
            # Convert to text representation
            text = df.to_string(index=False)
            
            row_count = len(df)
            
            return text.strip(), row_count
            
        except Exception as e:
            logger.error(f"CSV text extraction failed: {e}")
            raise Exception(f"CSV processing failed: {str(e)}")
    
    def process_file(self, file: UploadFile) -> ProcessedContent:
        """
        Process uploaded file and extract text content.
        
        Args:
            file: FastAPI UploadFile object
            
        Returns:
            ProcessedContent: Extracted text and metadata
            
        Raises:
            HTTPException: If file processing fails
        """
        import time
        start_time = time.time()
        
        try:
            # Validate file
            is_valid, error_msg = self.validate_file(file)
            if not is_valid:
                raise HTTPException(status_code=400, detail=error_msg)
            
            # Read file content
            content = file.file.read()
            file_size = len(content)
            
            # Calculate file hash
            file_hash = self.calculate_file_hash(content)
            
            # Determine file type
            file_ext = Path(file.filename).suffix.lower()
            mime_type = self.SUPPORTED_EXTENSIONS.get(file_ext, 'application/octet-stream')
            
            # Extract text based on file type
            text = ""
            page_count = None
            
            if file_ext == '.pdf':
                text, page_count = self.extract_text_from_pdf(content)
            elif file_ext in ['.docx', '.doc']:
                text, page_count = self.extract_text_from_docx(content)
            elif file_ext in ['.pptx', '.ppt']:
                text, page_count = self.extract_text_from_pptx(content)
            elif file_ext in ['.xlsx', '.xls']:
                text, page_count = self.extract_text_from_xlsx(content)
            elif file_ext == '.txt':
                text, page_count = self.extract_text_from_txt(content)
            elif file_ext in ['.html', '.htm']:
                text, page_count = self.extract_text_from_html(content)
            elif file_ext == '.md':
                text, page_count = self.extract_text_from_markdown(content)
            elif file_ext == '.json':
                text, page_count = self.extract_text_from_json(content)
            elif file_ext == '.csv':
                text, page_count = self.extract_text_from_csv(content)
            else:
                raise HTTPException(status_code=400, detail=f"Unsupported file type: {file_ext}")
            
            # Create metadata
            processing_time = time.time() - start_time
            metadata = FileMetadata(
                filename=file.filename,
                file_size=file_size,
                file_type=file_ext,
                mime_type=mime_type,
                file_hash=file_hash,
                processing_time=processing_time,
                text_length=len(text),
                page_count=page_count
            )
            
            # Create processed content
            processed_content = ProcessedContent(
                text=text,
                metadata=metadata
            )
            
            logger.info(f"Successfully processed file: {file.filename} ({file_size} bytes, {len(text)} chars)")
            
            return processed_content
            
        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"File processing failed for {file.filename}: {e}")
            processing_time = time.time() - start_time
            
            # Create error metadata
            metadata = FileMetadata(
                filename=file.filename,
                file_size=0,
                file_type=Path(file.filename).suffix.lower() if file.filename else "",
                mime_type="",
                file_hash="",
                processing_time=processing_time,
                text_length=0,
                error=str(e)
            )
            
            raise HTTPException(
                status_code=500, 
                detail=f"File processing failed: {str(e)}"
            )
    
    def process_multiple_files(self, files: List[UploadFile]) -> List[ProcessedContent]:
        """
        Process multiple files in batch.
        
        Args:
            files: List of FastAPI UploadFile objects
            
        Returns:
            List of ProcessedContent objects
        """
        results = []
        
        for file in files:
            try:
                processed_content = self.process_file(file)
                results.append(processed_content)
            except Exception as e:
                logger.error(f"Failed to process file {file.filename}: {e}")
                # Create error result
                metadata = FileMetadata(
                    filename=file.filename,
                    file_size=0,
                    file_type="",
                    mime_type="",
                    file_hash="",
                    processing_time=0,
                    text_length=0,
                    error=str(e)
                )
                results.append(ProcessedContent(
                    text="",
                    metadata=metadata
                ))
        
        return results
    
    def get_supported_formats(self) -> Dict[str, str]:
        """
        Get list of supported file formats.
        
        Returns:
            Dictionary mapping file extensions to MIME types
        """
        return self.SUPPORTED_EXTENSIONS.copy()
    
    def get_processing_stats(self) -> Dict[str, Any]:
        """
        Get file processing statistics.
        
        Returns:
            Dictionary with processing statistics
        """
        return {
            "supported_formats": len(self.SUPPORTED_EXTENSIONS),
            "max_file_size_mb": self.MAX_FILE_SIZE / (1024 * 1024),
            "upload_directory": str(self.upload_dir),
            "supported_extensions": list(self.SUPPORTED_EXTENSIONS.keys())
        }


# Global file processor instance
_file_processor: Optional[FileProcessor] = None


def get_file_processor() -> FileProcessor:
    """
    Get the global file processor instance.
    
    Returns:
        FileProcessor: Global file processor instance
    """
    global _file_processor
    if _file_processor is None:
        _file_processor = FileProcessor()
    return _file_processor


def initialize_file_processor(config_manager=None) -> FileProcessor:
    """
    Initialize the global file processor.
    
    Args:
        config_manager: Configuration manager instance
        
    Returns:
        FileProcessor: Initialized file processor instance
    """
    global _file_processor
    _file_processor = FileProcessor(config_manager)
    return _file_processor


def reset_file_processor() -> None:
    """Reset the global file processor instance (useful for testing)."""
    global _file_processor
    _file_processor = None
